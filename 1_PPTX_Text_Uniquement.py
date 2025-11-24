"""
VERSION 1 : TEXTE UNIQUEMENT
----------------------------
Ce script génère une présentation PowerPoint (.pptx) focalisée sur la structure et le contenu.
Nuance technique :
- Aucune dépendance externe (pas de requests, pas de torch).
- Utilisation de `python-pptx` pur.
- Création manuelle des zones de texte (TextBox) pour s'affranchir des layouts par défaut
  souvent rigides, garantissant un alignement précis au pixel près via l'unité 'Inches'.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def configurer_titre(slide, texte):
    # Création d'une boite de titre manuelle pour contrôle total
    titre_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    tf = titre_box.text_frame
    p = tf.paragraphs[0]
    p.text = texte
    p.font.name = 'Arial'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = None  # Couleur par défaut (Noir)


def ajouter_slide_texte(pres, titre, contenu_liste):
    # Layout 6 = Vide (Blank) pour tout construire soi-même
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    configurer_titre(slide, titre)

    # Zone de contenu principal
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
    tf = body_box.text_frame
    tf.word_wrap = True

    for item in contenu_liste:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = 'Arial'
        p.font.size = Pt(20)
        p.space_after = Pt(14)
        p.level = 0


def main():
    pres = Presentation()

    # Slide de Titre (Layout 0 standard)
    slide0 = pres.slides.add_slide(pres.slide_layouts[0])
    slide0.shapes.title.text = "Drones & Last Mile Delivery"
    slide0.placeholders[1].text = "Analyse Stratégique - Version Texte\nMichel - HEC Paris"

    # Données
    slides_data = [
        ("Contexte Économique", [
            "Coût du dernier km : 53% du coût total de la chaine logistique.",
            "Inefficacité des livraisons par camionnette en hyper-centre.",
            "Croissance e-commerce : +15% CAGR, saturation des infrastructures."
        ]),
        ("Architecture Technique", [
            "Vecteurs : Multirotors (court rayon) vs VTOL (long rayon).",
            "Navigation : Fusion de capteurs (LIDAR, GPS RTK, Flux optique).",
            "Énergie : Densité énergétique des Li-Ion vs Hydrogène."
        ]),
        ("Régulation & Avenir", [
            "Cadre EASA : SORA (Specific Operations Risk Assessment).",
            "U-Space : Gestion digitale du trafic basse altitude.",
            "Acceptabilité sociale : Nuisance sonore et vie privée."
        ])
    ]

    for titre, points in slides_data:
        ajouter_slide_texte(pres, titre, points)

    pres.save("1_Presentation_Texte.pptx")
    print("Génération terminée : 1_Presentation_Texte.pptx")


if __name__ == "__main__":
    main()