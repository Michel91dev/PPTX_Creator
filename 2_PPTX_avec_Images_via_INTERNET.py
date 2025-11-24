"""
VERSION 2 : IMAGES VIA INTERNET (WEB SCRAPING/API)
--------------------------------------------------
Ce script intègre des visuels récupérés dynamiquement sur le web.
Nuance technique :
- Utilisation de la librairie `requests` pour les appels HTTP GET.
- Gestion des flux binaires via `io.BytesIO` (l'image n'est jamais sauvegardée sur le disque,
  elle reste en RAM avant d'être injectée dans le XML du PPTX).
- Robustesse : Si l'image ne charge pas (timeout/404), le script ne plante pas.
"""

import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt


def telecharger_image(url_image):
    try:
        # Timeout court pour éviter de bloquer le script
        r = requests.get(url_image, timeout=5)
        if r.status_code == 200:
            return BytesIO(r.content)
    except Exception as e:
        print(f"Erreur téléchargement : {e}")
    return None


def ajouter_slide_web(pres, titre, points, url_img):
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # Titre
    titre_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    titre_box.text_frame.text = titre
    titre_box.text_frame.paragraphs[0].font.size = Pt(32)
    titre_box.text_frame.paragraphs[0].font.bold = True

    # Contenu (Gauche)
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = body_box.text_frame
    tf.word_wrap = True
    for point in points:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.space_after = Pt(12)

    # Image (Droite) téléchargée à la volée
    flux_img = telecharger_image(url_img)
    if flux_img:
        slide.shapes.add_picture(flux_img, Inches(5.2), Inches(1.8), width=Inches(4.5))
    else:
        # Fallback textuel si pas d'internet
        err_box = slide.shapes.add_textbox(Inches(5.2), Inches(3), Inches(4), Inches(1))
        err_box.text_frame.text = "[Image non disponible]"


def main():
    pres = Presentation()

    # Slide Titre
    s0 = pres.slides.add_slide(pres.slide_layouts[0])
    s0.shapes.title.text = "Drones Last Mile"
    s0.placeholders[1].text = "Version : Images Web\nMichel - HEC Paris"

    # Utilisation de placehold.co pour simuler des images parfaites pour la démo
    # Dans la réalité, remplacez l'URL par un lien direct vers un JPG/PNG
    slides = [
        {
            "titre": "Problématique Urbaine",
            "points": ["Congestion et pollution.", "Coûts logistiques explosifs."],
            "url": "https://placehold.co/800x600/png?text=Urban+Congestion"
        },
        {
            "titre": "Technologie Drone",
            "points": ["VTOL électrique.", "Navigation autonome AI."],
            "url": "https://placehold.co/800x600/png?text=Drone+Tech"
        }
    ]

    for s in slides:
        ajouter_slide_web(pres, s["titre"], s["points"], s["url"])

    pres.save("2_Presentation_Web.pptx")
    print("Génération terminée : 2_Presentation_Web.pptx")


if __name__ == "__main__":
    main()