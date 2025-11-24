"""
VERSION 3 : GÉNÉRATION IA LOCALE (APPLE SILICON / METAL)
--------------------------------------------------------
Ce script transforme le Mac en studio de création graphique autonome.
Nuance technique :
- Utilise `torch` et `diffusers` avec le backend 'mps' (Metal Performance Shaders).
- Optimisation float16 pour réduire l'usage RAM unifiée et accélérer le rendu.
- Modèle 'Stable Diffusion v1.5' chargé localement (4Go+ au premier run).
- Pipeline complet : Texte -> Tenseurs GPU -> Image PIL -> BytesIO -> PPTX.
"""

import torch
from diffusers import StableDiffusionPipeline
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

# --- Initialisation du moteur IA (Singleton) ---
print(">>> Chargement du modèle Stable Diffusion sur Mac Silicon (MPS)...")
try:
    MODEL_ID = "runwayml/stable-diffusion-v1-5"
    # float16 est crucial sur Mac pour la vitesse
    pipe = StableDiffusionPipeline.from_pretrained(MODEL_ID, torch_dtype=torch.float16)
    pipe = pipe.to("mps")
    pipe.enable_attention_slicing()  # Économie de mémoire
    print(">>> Moteur prêt.")
except Exception as e:
    print(f"ERREUR IA : {e}")
    pipe = None


def generer_et_inserer(pres, slide, prompt):
    if pipe is None:
        return

    print(f"   [IA] Génération : {prompt}...")
    # Inference (Calculs matriciels sur le GPU)
    image = pipe(prompt, num_inference_steps=30).images[0]

    # Conversion mémoire
    img_stream = BytesIO()
    image.save(img_stream, format="PNG")
    img_stream.seek(0)

    # Insertion dans la slide
    slide.shapes.add_picture(img_stream, Inches(5.2), Inches(1.8), width=Inches(4.5))


def ajouter_slide_ia(pres, titre, points, prompt_visuel):
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # Titre
    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    t_box.text_frame.text = titre
    t_box.text_frame.paragraphs[0].font.size = Pt(32)
    t_box.text_frame.paragraphs[0].font.bold = True

    # Texte
    b_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = b_box.text_frame
    tf.word_wrap = True
    for p_txt in points:
        p = tf.add_paragraph()
        p.text = f"• {p_txt}"
        p.font.size = Pt(16)
        p.space_after = Pt(12)

    # Appel de la fonction IA
    generer_et_inserer(pres, slide, prompt_visuel)


def main():
    pres = Presentation()

    s0 = pres.slides.add_slide(pres.slide_layouts[0])
    s0.shapes.title.text = "Drones Last Mile"
    s0.placeholders[1].text = "Version : IA Générative Locale (MPS)\nMichel - HEC Paris"

    data = [
        {
            "titre": "Livraison Urbaine",
            "points": ["Rapidité extrême.", "Zéro émission locale."],
            "prompt": "Futuristic delivery drone flying in a city street, photorealistic, 4k, cinematic lighting"
        },
        {
            "titre": "Contrôle à distance",
            "points": ["Pilotage BVLOS.", "Télémetrie 5G."],
            "prompt": "Operator looking at screens controlling drones, cyberpunk style, digital art, detailed"
        }
    ]

    for d in data:
        ajouter_slide_ia(pres, d["titre"], d["points"], d["prompt"])

    pres.save("3_Presentation_IA_Local.pptx")
    print("Génération terminée : 3_Presentation_IA_Local.pptx")


if __name__ == "__main__":
    main()