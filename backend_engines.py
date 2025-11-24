import re
import requests
import torch
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from diffusers import StableDiffusionPipeline

# --- UTILITAIRES ---
def init_presentation(titre_doc, sous_titre):
    """Crée la base de la présentation"""
    pres = Presentation()
    slide0 = pres.slides.add_slide(pres.slide_layouts[0])
    slide0.shapes.title.text = titre_doc
    slide0.placeholders[1].text = sous_titre
    return pres

def add_slide_layout(pres, titre, points, image_stream=None):
    """Ajoute une slide standardisée (Titre + Liste + Image optionnelle)"""
    layout_vide = pres.slide_layouts[6]
    slide = pres.slides.add_slide(layout_vide)

    # Titre
    box_titre = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    tf = box_titre.text_frame
    tf.text = titre
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = 'Arial'

    # Texte : largeur différente selon la présence d'une image
    if image_stream:
        # Mode texte + image : colonne gauche plus étroite
        box_txt = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(5))
    else:
        # Mode texte seul : zone de texte élargie sur la droite
        box_txt = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf_txt = box_txt.text_frame
    tf_txt.word_wrap = True
    for point in points:
        p = tf_txt.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.space_after = Pt(14)
        p.font.name = 'Arial'

    # Image (Droite)
    if image_stream:
        slide.shapes.add_picture(image_stream, Inches(5.8), Inches(1.8), width=Inches(3.8))

# --- MOTEUR 1 : TEXTE SEUL ---
def generate_text_only(data_slides, filename="Sortie_Texte.pptx", image_files=None):
    """Génère une présentation texte seul.

    image_files : liste optionnelle de fichiers Streamlit uploadés.
    - Les images sont appliquées en fonction du numéro au début du nom du fichier (1 à 99).
      Exemple : "01 campus.jpg" ou "1 campus.jpg" -> slide 1 ; "10 salle.png" -> slide 10.
    - Si aucune image n'est trouvée pour un numéro de slide donné, la slide reste en texte seul.
    """
    pres = init_presentation("Présentation Texte", "Mode Rapide - HEC")

    image_files = image_files or []

    # Construction d'un mapping numero_de_slide -> contenu d'image
    # On lit 1 ou 2 chiffres au tout début du nom de fichier : 1..99
    image_map = {}
    for f in image_files:
        try:
            name = getattr(f, "name", "") or ""
            m = re.match(r"^(\d{1,2})", name.strip())
            if not m:
                continue
            numero = int(m.group(1))
            if 1 <= numero <= 99 and numero not in image_map:
                img_bytes = f.read()
                image_map[numero] = img_bytes
        except Exception as e:
            print(f"Erreur préparation image '{getattr(f, 'name', '?')}' : {e}")

    for idx, s in enumerate(data_slides):
        titre_brut = s['titre']
        image_stream = None

        # Numéro de slide logique (1-based)
        slide_num = idx + 1
        img_bytes = image_map.get(slide_num)
        if img_bytes is not None:
            try:
                image_stream = BytesIO(img_bytes)
            except Exception as e:
                print(f"Erreur lecture image pour le slide {slide_num}: {e}")

        add_slide_layout(pres, titre_brut, s['points'], image_stream=image_stream)
    
    buffer = BytesIO()
    pres.save(buffer)
    buffer.seek(0)
    return buffer

# --- MOTEUR 2 : WEB IMAGES ---
def generate_web_images(data_slides, filename="Sortie_Web.pptx"):
    pres = init_presentation("Présentation Web", "Mode Connecté - HEC")
    
    for s in data_slides:
        url = s.get('visuel', '') # Ici 'visuel' contient l'URL
        img_stream = None
        if url.startswith("http"):
            try:
                r = requests.get(url, timeout=4)
                if r.status_code == 200:
                    img_stream = BytesIO(r.content)
            except:
                pass
        
        add_slide_layout(pres, s['titre'], s['points'], img_stream)

    buffer = BytesIO()
    pres.save(buffer)
    buffer.seek(0)
    return buffer

# --- MOTEUR 3 : LOCAL AI (MAC SILICON) ---
# On prépare le chargement du modèle mais on ne l'exécute que si nécessaire
_pipe_cache = None

def get_ai_pipeline():
    global _pipe_cache
    if _pipe_cache is not None:
        return _pipe_cache
    
    model_id = "runwayml/stable-diffusion-v1-5"

    # 1) Tentative MPS float16 (préférée sur Apple Silicon)
    if torch.backends.mps.is_available():
        try:
            print(">>> Chargement Modèle IA (MPS, float16)...")
            pipe = StableDiffusionPipeline.from_pretrained(model_id, torch_dtype=torch.float16)
            pipe = pipe.to("mps")

            pipe.enable_attention_slicing()

            if hasattr(pipe, "safety_checker"):
                pipe.safety_checker = None

            _pipe_cache = pipe
            return pipe
        except Exception as e:
            print(f"Erreur chargement IA sur MPS, bascule sur CPU: {e}")

    # 2) Fallback CPU float32
    try:
        print(">>> Chargement Modèle IA (CPU, float32)...")
        pipe = StableDiffusionPipeline.from_pretrained(model_id, torch_dtype=torch.float32)
        pipe = pipe.to("cpu")

        pipe.enable_attention_slicing()

        if hasattr(pipe, "safety_checker"):
            pipe.safety_checker = None

        _pipe_cache = pipe
        return pipe
    except Exception as e:
        print(f"Erreur chargement IA (CPU): {e}")
        return None

def generate_local_ai(data_slides, progress_callback=None, num_steps=30, per_slide_images=False):
    pipe = get_ai_pipeline()
    if not pipe:
        return None

    pres = init_presentation("Présentation AI", "Génération Locale Mac Silicon")
    
    total = len(data_slides)
    for i, s in enumerate(data_slides):
        prompt = s.get('visuel', '') # Ici 'visuel' contient le prompt
        
        # Mise à jour barre de progression UI
        if progress_callback:
            progress_callback(i / total, f"Génération visuel {i+1}/{total} : {prompt[:30]}...")

        img_stream = None

        # Génération IA : soit uniquement pour la première slide, soit pour chaque slide
        doit_generer = False
        if per_slide_images and prompt:
            doit_generer = True
        elif (not per_slide_images) and i == 0 and prompt:
            doit_generer = True

        if doit_generer:
            try:
                image = pipe(prompt, num_inference_steps=num_steps).images[0]
                img_stream = BytesIO()
                image.save(img_stream, format="PNG")
                img_stream.seek(0)
            except Exception as e:
                print(f"Erreur génération image IA pour la slide {i+1}: {e}")

        add_slide_layout(pres, s['titre'], s['points'], img_stream)

    buffer = BytesIO()
    pres.save(buffer)
    buffer.seek(0)
    return buffer